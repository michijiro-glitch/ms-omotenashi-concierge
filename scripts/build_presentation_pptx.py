# -*- coding: utf-8 -*-
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
PUB = ROOT / "public"
OUT = ROOT / "卒業発表_MsOmotenashiConcierge.pptx"

IVORY = RGBColor(0xFB, 0xF8, 0xF2)
NAVY = RGBColor(0x24, 0x30, 0x56)
GOLD = RGBColor(0xC4, 0xA3, 0x6A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT = RGBColor(0x6B, 0x5E, 0x4E)
LINE = RGBColor(0xD7, 0xCC, 0xBA)

W = Inches(13.333)
H = Inches(7.5)
SERIF = "Yu Mincho"
SANS = "Yu Gothic"


def rgb_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def line_color(shape, color, width=Pt(1)):
    shape.line.color.rgb = color
    shape.line.width = width


def set_run(run, text, size, color, bold=False, font=SANS):
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font
    rPr = run._r.get_or_add_rPr()
    latin = rPr.find(qn("a:latin"))
    if latin is None:
        latin = rPr.makeelement(qn("a:latin"), {"typeface": font})
        rPr.append(latin)
    else:
        latin.set("typeface", font)
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {"typeface": font})
        rPr.append(ea)
    else:
        ea.set("typeface", font)


def textbox(slide, left, top, width, height, lines, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, spec in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = spec.get("align", align)
        p.space_after = Pt(spec.get("after", 6))
        run = p.add_run()
        set_run(
            run,
            spec["text"],
            spec.get("size", 14),
            spec.get("color", NAVY),
            spec.get("bold", False),
            spec.get("font", SANS),
        )
    return box


def kicker(slide, text):
    textbox(
        slide,
        Inches(0.7),
        Inches(0.28),
        Inches(12),
        Inches(0.35),
        [{"text": text, "size": 11, "color": GOLD, "font": SANS, "after": 0, "align": PP_ALIGN.CENTER}],
        align=PP_ALIGN.CENTER,
    )


def title(slide, text, top=0.55):
    textbox(
        slide,
        Inches(0.7),
        Inches(top),
        Inches(12),
        Inches(0.55),
        [{"text": text, "size": 26, "color": NAVY, "font": SERIF, "after": 0, "align": PP_ALIGN.CENTER}],
        align=PP_ALIGN.CENTER,
    )


def gold_rule(slide, top=1.15):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.4), Inches(top), Inches(2.5), Pt(1.2))
    rgb_fill(shape, GOLD)
    shape.line.fill.background()
    dia = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(6.5), Inches(top) - Pt(4), Pt(10), Pt(10))
    rgb_fill(dia, IVORY)
    line_color(dia, GOLD, Pt(1))


def card(slide, left, top, width, height, heading, body):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    rgb_fill(shape, WHITE)
    line_color(shape, LINE, Pt(1))
    try:
        shape.adjustments[0] = 0.08
    except Exception:
        pass
    textbox(
        slide,
        left + Inches(0.16),
        top + Inches(0.14),
        width - Inches(0.32),
        Inches(0.36),
        [{"text": heading, "size": 14, "color": NAVY, "font": SERIF, "after": 0}],
    )
    body_lines = [{"text": line, "size": 13, "color": SOFT, "font": SANS, "after": 8} for line in body.split("\n") if line]
    if body_lines:
        body_lines[-1]["after"] = 0
    textbox(
        slide,
        left + Inches(0.16),
        top + Inches(0.5),
        width - Inches(0.32),
        height - Inches(0.62),
        body_lines,
    )


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def add_fitted_picture(slide, path, left, top, width, height, mode="cover"):
    """Fit a screenshot into a box. cover = fill from the top; contain = letterbox."""
    pic = slide.shapes.add_picture(str(path), left, top)
    native_w, native_h = pic.width, pic.height
    ratio = native_w / native_h
    box_ratio = width / height
    if mode == "contain":
        if ratio > box_ratio:
            pic.width = width
            pic.height = int(width / ratio)
        else:
            pic.height = height
            pic.width = int(height * ratio)
        pic.left = int(left + (width - pic.width) / 2)
        pic.top = int(top + (height - pic.height) / 2)
        return pic
    if ratio > box_ratio:
        pic.height = height
        pic.width = int(height * ratio)
        overflow = pic.width - width
        crop = overflow / pic.width / 2
        pic.crop_left = crop
        pic.crop_right = crop
        pic.width = width
        pic.left = left
        pic.top = top
    else:
        pic.width = width
        pic.height = int(width / ratio)
        extra = 1 - (height / pic.height)
        pic.crop_bottom = extra
        pic.height = height
        pic.left = left
        pic.top = top
    return pic


def add_picture_box(slide, path, left, top, width, height, caption, mode="cover"):
    frame = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rgb_fill(frame, WHITE)
    line_color(frame, LINE, Pt(1))
    cap_h = Inches(0.28)
    pad = Inches(0.04)
    img_top = top + pad
    img_left = left + pad
    img_w = width - pad * 2
    img_h = height - cap_h - pad
    if Path(path).exists():
        add_fitted_picture(slide, path, img_left, img_top, img_w, img_h, mode=mode)
    textbox(
        slide,
        left + Inches(0.08),
        top + height - cap_h,
        width - Inches(0.16),
        cap_h,
        [{"text": caption, "size": 10, "color": SOFT, "font": SANS, "after": 0}],
    )


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    blank_layout = prs.slide_layouts[6]

    def new():
        s = prs.slides.add_slide(blank_layout)
        fill = s.background.fill
        fill.solid()
        fill.fore_color.rgb = IVORY
        return s

    s = new()
    kicker(s, "ADS 卒業制作｜開発ストーリー")
    logo = PUB / "logo.png"
    if logo.exists():
        s.shapes.add_picture(str(logo), Inches(5.15), Inches(1.15), width=Inches(3.0))
    gold_rule(s, 3.05)
    textbox(
        s,
        Inches(1.2),
        Inches(3.25),
        Inches(10.9),
        Inches(1.1),
        [
            {"text": "出版社の現場で培った目利きで選ぶ、", "size": 20, "color": NAVY, "font": SERIF, "after": 4, "align": PP_ALIGN.CENTER},
            {"text": "レストランとギフトのパーソナルガイド", "size": 20, "color": NAVY, "font": SERIF, "after": 0, "align": PP_ALIGN.CENTER},
        ],
        align=PP_ALIGN.CENTER,
    )
    textbox(
        s,
        Inches(1.6),
        Inches(4.55),
        Inches(10.1),
        Inches(2.2),
        [
            {"text": "頭の中の経験を、再利用できる知識にする。", "size": 14, "color": SOFT, "after": 6, "align": PP_ALIGN.CENTER},
            {"text": "口コミサイトではなく、M自身の選択眼を残すパーソナル・コンシェルジュ。", "size": 14, "color": SOFT, "after": 6, "align": PP_ALIGN.CENTER},
            {"text": "広告営業出身。エンジニアではない。実装は Cursor。主役は、何を知識として残すかという設計。", "size": 13, "color": SOFT, "after": 0, "align": PP_ALIGN.CENTER},
        ],
        align=PP_ALIGN.CENTER,
    )
    notes(
        s,
        "表紙。すぐスライド2へ。\n"
        "結論：グルメアプリではなく、自分の目利きで選んだ経験を、検索できる知識に変えた。"
        "AIは後から付けるために、先に構造化した。",
    )

    s = new()
    kicker(s, "解決する面倒・課題")
    title(s, "思い出すだけでなく、呼び出せる形に")
    gold_rule(s)
    q1 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.45), Inches(5.7), Inches(1.15))
    q2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.45), Inches(5.7), Inches(1.15))
    for q in (q1, q2):
        rgb_fill(q, WHITE)
        line_color(q, LINE, Pt(1))
        try:
            q.adjustments[0] = 0.08
        except Exception:
            pass
    textbox(s, Inches(0.95), Inches(1.7), Inches(5.4), Inches(0.7), [{"text": "「喜んでもらえるかな」", "size": 20, "color": NAVY, "font": SERIF, "after": 0}], align=PP_ALIGN.CENTER)
    textbox(s, Inches(6.95), Inches(1.7), Inches(5.4), Inches(0.7), [{"text": "「センスのいいひとだと思われたい」", "size": 18, "color": NAVY, "font": SERIF, "after": 0}], align=PP_ALIGN.CENTER)
    card(s, Inches(0.7), Inches(2.85), Inches(3.85), Inches(2.55), "会食の店選び", "クライアントの会食。気の置けない代理店やベンダーとの大人数。「どこにする？」。苦手な食材がある人もいる。")
    card(s, Inches(4.75), Inches(2.85), Inches(3.85), Inches(2.55), "手土産", "いつもワンパターンの菓子折り。センスを見せたい。")
    card(s, Inches(8.8), Inches(2.85), Inches(3.85), Inches(2.55), "お祝いの花", "昇格のお祝いなど。センスの良いものにしたい。")
    textbox(
        s,
        Inches(0.8),
        Inches(5.55),
        Inches(11.7),
        Inches(1.4),
        [
            {"text": "お店選びも、ギフトの品選びも、その人を表す。", "size": 15, "color": NAVY, "font": SERIF, "after": 6, "align": PP_ALIGN.CENTER},
            {"text": "判断基準は頭の中にあった。でも、その場で思い出すしかなかった。", "size": 14, "color": SOFT, "after": 0, "align": PP_ALIGN.CENTER},
        ],
        align=PP_ALIGN.CENTER,
    )
    notes(
        s,
        "クライアントとの会食が決まって、店を選ぶ。手土産も選ぶ。\n"
        "いつも一緒に仕事をしている代理店やベンダーと、気の置けない会食。大人数だけど、どこにする？\n"
        "苦手な食材がある人もいる。お土産は、いつもワンパターンの菓子折り。でも、センスを見せたい。\n"
        "昇格のお祝いには、花。センスの良いものにしたい。\n\n"
        "心の中では、「喜んでもらえるかな」。そして、「センスのいいひとだと思われたい」。\n"
        "お店選びも、ギフトの品選びも、その人を表すものだと、思っています。\n"
        "判断基準は、頭の中にはあった。でも、その場で思い出すしかなかったんです。\n\n"
        "結論：グルメアプリではなく、自分の目利きで選んだ経験を、検索できる知識に変えた。"
        "AIは後から付けるために、先に構造化した。",
    )

    s = new()
    kicker(s, "作りたいと思った理由")
    title(s, "実用的で、記録になり、育てられる")
    gold_rule(s)
    card(s, Inches(0.7), Inches(1.5), Inches(3.85), Inches(3.4), "実用的", "相談されたその場で、条件から店や手土産を呼び出せる。")
    card(s, Inches(4.75), Inches(1.5), Inches(3.85), Inches(3.4), "自分の記録にもなる", "行った店、贈りたい手土産を、自分用の備忘録として残せる。")
    card(s, Inches(8.8), Inches(1.5), Inches(3.85), Inches(3.4), "これからもアップデートできる", "シートに足せば公開側も育つ。卒業制作で止めない前提。")
    textbox(
        s,
        Inches(0.8),
        Inches(5.2),
        Inches(11.7),
        Inches(1.4),
        [
            {"text": "口コミの量ではなく、「誰が、どんな基準で選ぶか」。", "size": 15, "color": NAVY, "font": SERIF, "after": 6, "align": PP_ALIGN.CENTER},
            {"text": "Restaurant と Gift に絞る。", "size": 14, "color": SOFT, "after": 0, "align": PP_ALIGN.CENTER},
        ],
        align=PP_ALIGN.CENTER,
    )
    notes(
        s,
        "お題を選んだのは、実用的で、自分の記録にもなるし、これからもアップデートできるからです。"
        "相談されたその場で呼び出せる。行った店や贈りたい手土産を、備忘録として残せる。"
        "シートに足せば、公開側も育つ。卒業制作で止めない前提です。",
    )

    s = new()
    kicker(s, "作ったもの")
    title(s, "見る・探すは公開。直すのは自分用")
    gold_rule(s, 1.05)
    shots = [
        (PUB / "images/home/_verify-home.png", "公開｜トップ　Restaurant / Gift"),
        (PUB / "images/home/_verify-restaurants.png", "公開｜Restaurant　条件・単語・音声"),
        (PUB / "images/about/_preview/about-visit-wishlist.png", "公開｜Visit / Wishlist　店舗ステータス"),
        (PUB / "images/edit/_verify-edit.png", "入力｜非公開 /edit　フォームから移した"),
    ]
    # 1400x900 を切らずに収める。横長枠で cover すると Restaurant / Gift が消える。
    box_w = Inches(4.55)
    box_h = Inches(3.12)
    gap_x, gap_y = Inches(0.16), Inches(0.1)
    grid_w = box_w * 2 + gap_x
    left0 = (W - grid_w) / 2
    top0 = Inches(1.18)
    positions = [
        (left0, top0),
        (left0 + box_w + gap_x, top0),
        (left0, top0 + box_h + gap_y),
        (left0 + box_w + gap_x, top0 + box_h + gap_y),
    ]
    for (path, cap), (l, t) in zip(shots, positions):
        add_picture_box(s, path, l, t, box_w, box_h, cap, mode="contain")
    notes(
        s,
        "公開は、見る・探すだけ。直すのは自分用。最初は Google フォーム、いまは非公開の画面から入力します。"
        "アプリで直しても、シートに残る。",
    )

    s = new()
    kicker(s, "作ったもの｜About M")
    title(s, "誰が、どんな基準で選ぶか")
    gold_rule(s)
    about = PUB / "images/about/_preview/about-page-desktop.png"
    add_picture_box(
        s,
        about,
        Inches(1.9),
        Inches(1.4),
        Inches(9.5),
        Inches(5.35),
        "公開｜About M　会食・ワイン・贈答。「美味しいだけでなく、その場にちょうどいい」",
        mode="contain",
    )
    notes(
        s,
        "口コミの量ではなく、私自身の選択眼を残す、パーソナル・コンシェルジュです。\n"
        "出版社の現場で培った目利き。「美味しいだけでなく、その場にちょうどいい」。"
        "「誰が、どんな基準で選ぶか」。それを置くのが、About M です。\n"
        "デモでも、このページを開きます。",
    )

    s = new()
    kicker(s, "工夫したポイントと苦戦したポイント")
    title(s, "経験を、検索できる項目にした")
    gold_rule(s)
    card(s, Inches(0.7), Inches(1.5), Inches(5.85), Inches(2.15), "データ化", "メモ保存ではなく、呼び出せる項目にする。Visit＝実際に訪問。Wishlist＝行ってみたい。Choiceは特集で、ステータスではない。")
    card(s, Inches(6.75), Inches(1.5), Inches(5.85), Inches(2.15), "About M", "口コミの量では勝負しない。「誰が、どんな経験と基準で選ぶか」。美味しいだけでなく、その場にちょうどいい。")
    card(s, Inches(0.7), Inches(3.85), Inches(5.85), Inches(2.15), "入力を非公開へ", "当初は Google フォーム。いまは非公開の /edit。来客向けから登録・直すを外した。アプリで直しても、シートに残る。")
    card(s, Inches(6.75), Inches(3.85), Inches(5.85), Inches(2.15), "デザインと音声", "見た目は Cursor。かわいい案から上質・知的・ニュートラルへ。いまの音声は入力。「恵比寿 和食」はできる。文の意味理解はこれから。")
    textbox(
        s,
        Inches(0.8),
        Inches(6.2),
        Inches(11.7),
        Inches(0.7),
        [{"text": "AIを付けることから始めたのではない。先に、自分の経験を構造化した。", "size": 15, "color": NAVY, "font": SERIF, "after": 0, "align": PP_ALIGN.CENTER}],
        align=PP_ALIGN.CENTER,
    )
    notes(
        s,
        "やり方は、単純です。暗黙知を、検索できる項目にする。\n"
        "Visit は実際に訪問、Wishlist は行ってみたい。Choice は特集で、ステータスではありません。\n"
        "見た目も、AI で作りました。かわいい案から、上質・知的・ニュートラルへ。デザインは Cursor。検索の AI API は、まだ繋いでいません。\n"
        "いまの音声は、入力です。「恵比寿 和食」はできる。「恵比寿で和食が食べたい」の意味理解は、これから。"
        "世の中の店を AI に聞くのではありません。私の登録データを、再利用する。",
    )

    s = new()
    kicker(s, "デモ")
    title(s, "ここから、実演します", top=2.2)
    gold_rule(s, 2.85)
    textbox(
        s,
        Inches(0.7),
        Inches(3.2),
        Inches(12),
        Inches(1.2),
        [{"text": "トップ　→　About M　→　Restaurant　→　音声検索　→　1店舗", "size": 24, "color": NAVY, "font": SERIF, "after": 0, "align": PP_ALIGN.CENTER}],
        align=PP_ALIGN.CENTER,
    )
    textbox(
        s,
        Inches(1.5),
        Inches(4.6),
        Inches(10.3),
        Inches(1.2),
        [{"text": "Gift は操作しない。見た目のAIと、店を探すAIは分ける。", "size": 16, "color": SOFT, "after": 0, "align": PP_ALIGN.CENTER}],
        align=PP_ALIGN.CENTER,
    )
    notes(
        s,
        "いちばん見てほしいのは、ここです。\n"
        "トップ。About M。誰が選ぶかを、ページで見ます。\n"
        "Restaurant。音声で、単語を入れる。結果。1店舗を見る。\n"
        "Gift は操作しません。以上が、いま動いているものです。",
    )

    s = new()
    kicker(s, "コンセプトの続きと、これから")
    title(s, "検索から、知識へ")
    gold_rule(s)
    card(
        s,
        Inches(0.45),
        Inches(1.4),
        Inches(4.0),
        Inches(4.55),
        "現在",
        "・条件検索・音声検索\n・無料で利用者とデータを増やす\n・予約・購入へ送客する土台",
    )
    card(
        s,
        Inches(4.65),
        Inches(1.4),
        Inches(4.0),
        Inches(4.55),
        "次",
        "・AI自然文検索・登録支援\n・M's Choice 特集\n・noteで読む\n・アプリで探す・比較\n・予約／ECへつなぐ",
    )
    card(
        s,
        Inches(8.85),
        Inches(1.4),
        Inches(4.0),
        Inches(4.55),
        "将来",
        "・個別コンシェルジュ\n・法人の会食・贈答ナレッジ\n・この相手ならこの店\n・この場面ならこの手土産",
    )
    textbox(
        s,
        Inches(0.6),
        Inches(6.15),
        Inches(12.1),
        Inches(0.7),
        [{"text": "経験 → 構造化 → 検索 → 音声 → 公開 → AI → 個人 → 法人", "size": 14, "color": SOFT, "after": 0, "align": PP_ALIGN.CENTER}],
        align=PP_ALIGN.CENTER,
    )
    notes(
        s,
        "いまは、条件検索と音声検索。無料で利用者とデータを増やす。予約・購入へ送客する土台。\n"
        "次は、自然文を AI が解釈して、条件に変え、M の登録データから返す。"
        "M's Choice。詳しく読むのは note、探す・比較するのはこのアプリ、予約・購入は外。\n"
        "将来は、個別コンシェルジュと、法人の会食・贈答ナレッジ。"
        "この相手ならこの店、この場面ならこの手土産。",
    )

    s = new()
    kicker(s, "これから")
    title(s, "備忘録から、知識へ")
    gold_rule(s)
    textbox(
        s,
        Inches(0.8),
        Inches(2.0),
        Inches(11.7),
        Inches(4.4),
        [
            {"text": "資産はUIやコードではなく、経験・選定データ・蓄積された知識。", "size": 16, "color": SOFT, "after": 18, "align": PP_ALIGN.CENTER},
            {"text": "自分のための備忘録　→　シートとアプリ　→　暗黙知を再利用可能な知識へ", "size": 20, "color": NAVY, "font": SERIF, "after": 20, "align": PP_ALIGN.CENTER},
            {"text": "ADSのチームgraceの皆様、半年間ご一緒させていただきありがとうございました。", "size": 15, "color": NAVY, "after": 6, "align": PP_ALIGN.CENTER},
            {"text": "これからも学びを続けていきたいと思います。", "size": 15, "color": NAVY, "after": 16, "align": PP_ALIGN.CENTER},
            {"text": "是非見ていただいて、お気づきの点やアイデアなどいただけましたら嬉しいです。", "size": 15, "color": NAVY, "after": 8, "align": PP_ALIGN.CENTER},
            {"text": "マコなり社長ともぜひコラボしたいです。よろしくお願いします。", "size": 15, "color": NAVY, "after": 0, "align": PP_ALIGN.CENTER},
        ],
        align=PP_ALIGN.CENTER,
    )
    notes(
        s,
        "資産は、UI やコードではありません。経験、選定データ、蓄積された知識です。\n"
        "自分のための備忘録が、シートとアプリを経て、個人の暗黙知を再利用可能な知識に変わる。その先に、法人のナレッジがある。\n\n"
        "ADSのチームgraceの皆様、半年間ご一緒させていただきありがとうございました。"
        "これからも学びを続けていきたいと思います。\n"
        "是非見ていただいて、お気づきの点やアイデアなどいただけましたら嬉しいです。"
        "マコなり社長ともぜひコラボしたいです。よろしくお願いします。",
    )

    prs.save(OUT)
    desktop_dir = Path.home() / "OneDrive" / "デスクトップ"
    copied = None
    if desktop_dir.is_dir():
        dest = desktop_dir / "卒業発表.pptx"
        try:
            dest.write_bytes(OUT.read_bytes())
            copied = dest
        except OSError:
            dest = desktop_dir / "卒業発表_最新.pptx"
            dest.write_bytes(OUT.read_bytes())
            copied = dest
    print(str(OUT))
    if copied:
        print(str(copied))


if __name__ == "__main__":
    build()
